#!/usr/bin/env python
# pip install python-pptx pandas

import os
import csv
from pptx import Presentation
import pandas as pd
import argparse


# Set up argument parser
parser = argparse.ArgumentParser(description="Export speaker notes from a PowerPoint presentation to a CSV file.")
parser.add_argument('-i', '--input', required=True, help="Path to the input PowerPoint (.pptx) file")
parser.add_argument('-o', '--output', required=True, help="Path to the output CSV file")
parser.add_argument('--botnoi', action='store_true', help="Specify the output format to Botnoi")


# -----------------------------------------------------------------------------
def export_notes_to_csv(pptx_file, output_csv, args):
    # Load the PowerPoint presentation
    presentation = Presentation(pptx_file)

    # Create a list to hold the slide number and notes
    notes_data = []

    # Iterate through each slide
    for slide_number, slide in enumerate(presentation.slides, start=1):
        # Check if the slide has any notes
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text
            # convert multiline text to single line
            notes_text = notes_text.replace('\n', '  ').replace('\r', '')
        else:
            notes_text = ''

        # Append slide number and notes to list
        if args.botnoi:
            notes_data.append({'text': notes_text if notes_text else '.' , 'speaker_name': 'Ava'})
        else:
            notes_data.append({"Slide Number": f"[{slide_number}]", "Notes": notes_text})

    # Convert the data to a DataFrame
    df = pd.DataFrame(notes_data)

    # Export the DataFrame to CSV
    df.to_csv(output_csv, index=False, encoding='utf-8')
    print(f"Notes exported to {output_csv}")


# -----------------------------------------------------------------------------
if __name__ == '__main__':
    # Parse the arguments
    args = parser.parse_args()

    # Set the file paths from the parsed arguments
    pptx_file = args.input
    output_csv = args.output

    # Run the export function
    export_notes_to_csv(pptx_file, output_csv, args)
